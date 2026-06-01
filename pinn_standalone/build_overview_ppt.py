"""Generate overview PPTX with image placeholders for manual insertion."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "output" / "pinn_standalone_overview.pptx"


def add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def add_placeholder_slide(prs, title: str, hint: str, caption: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.75))
    tf = tx.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    left = Inches(0.7)
    top = Inches(1.35)
    width = Inches(9.1)
    height = Inches(4.35)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box.line.color.rgb = RGBColor(180, 180, 180)
    box.line.width = Pt(1.5)

    inner = slide.shapes.add_textbox(left, top + height / 2 - Inches(0.45), width, Inches(0.9))
    itf = inner.text_frame
    itf.paragraphs[0].text = "【插入图片】"
    itf.paragraphs[0].alignment = PP_ALIGN.CENTER
    itf.paragraphs[0].font.size = Pt(22)
    itf.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
    itf.vertical_anchor = MSO_ANCHOR.MIDDLE

    hint_box = slide.shapes.add_textbox(left, top + Inches(0.25), width, Inches(0.55))
    htf = hint_box.text_frame
    htf.paragraphs[0].text = hint
    htf.paragraphs[0].alignment = PP_ALIGN.CENTER
    htf.paragraphs[0].font.size = Pt(14)
    htf.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)

    cap = slide.shapes.add_textbox(Inches(0.7), top + height + Inches(0.15), Inches(9.1), Inches(0.85))
    ctf = cap.text_frame
    ctf.word_wrap = True
    ctf.paragraphs[0].text = caption
    ctf.paragraphs[0].font.size = Pt(14)


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "过氧化氢灭菌腔 · 多相多组分 PINN 代理模型",
        "固定网格 · 参数化入口 · 物理约束 + 数据驱动\n（概述稿 · 灰色框内可自行替换为示意图 / 截图）",
    )

    add_bullet_slide(
        prs,
        "整体概述",
        [
            "目标：在固定网格下，用 CFD 快照训练可泛化的时空场代理",
            "输入：坐标 + 离散时间 + T_preheat / T_h2o2（工况文件夹 + UDF）",
            "输出：20 维场（三相速度、p、T、k、ω、VOF、组分等）",
            "几何：外流体 fluid_o + 内流体 fluid_i + 固体壁 soild",
            "训练：数据拟合 + PDE 残差（无量纲）+ 结构约束 / 硬入口 + 软守恒",
        ],
    )

    add_placeholder_slide(
        prs,
        "图 1 · 网格分区空间示意",
        "建议：inspect_zones.py 输出的 zones_3d.png，或 Fluent 网格截图",
        "说明：内外两腔流体与固体壁面的空间关系；可在备注中标注 fluid_o / fluid_i / soild。",
    )

    add_placeholder_slide(
        prs,
        "图 2 · 模型输入输出（示意图）",
        "建议：方框图画 (x,y,z,t, bc) → MLP → 20 维输出",
        "说明：突出「工况参数 + 离散时间」作为条件输入；输出列表可与 FIELD_NAMES 对齐。",
    )

    add_placeholder_slide(
        prs,
        "图 3 · 物理约束与损失结构",
        "建议：流程图 — 数据损失 / PDE 各项 / 固体温度 / 入口 / 初值 / 组分·VOF 软约束",
        "说明：无量纲残差权重与 enabled_pdes 可在备注中列一行。",
    )

    add_placeholder_slide(
        prs,
        "图 4 · 数据与采样管线",
        "建议：cas.h5 + dat.h5 → 归一化 → 工况字典 → DataLoader → train_step",
        "说明：可标注 classify_fluid_cells、界面面心缓存、按工况划分 train/val。",
    )

    add_placeholder_slide(
        prs,
        "图 5 · 训练曲线或验证误差",
        "建议：Total loss / 按场的 MSE / Validation 曲线截图",
        "说明：汇报时选 1～2 条最能说明收敛与泛化的曲线即可。",
    )

    add_placeholder_slide(
        prs,
        "图 6 · 固体（或界面）温度可视化",
        "建议：solid_temp_viz 的 PNG，或 ParaView 打开的 VTK 截图",
        "说明：对比 Actual / Predicted / Error；可附 RMSE 文字在页脚自行添加。",
    )

    add_bullet_slide(
        prs,
        "固体壁面物性（当前配置）",
        [
            "ρ_s = 1000 kg/m³",
            "c_p,s = 1500 J/(kg·K)",
            "λ_s = 0.114 W/(m·K)（非金属壁，导热相对偏弱）",
        ],
    )

    add_bullet_slide(
        prs,
        "汇报提示（可选）",
        [
            "替换图片：选中灰色框 → 右键 → 更改图片；或删除框后插入图片并拉大至相近尺寸",
            "版式：宽屏 16:9 可在「设计」中改为 16:9（当前约 10×7.5 英寸）",
            "飞书文档：物理体系版 / 代码嵌入版 / PPT 底稿 三份相互对照",
        ],
    )

    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
